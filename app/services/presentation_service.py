# app/services/presentation_service.py
import io
from datetime import date

try:
    import matplotlib

    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
except ImportError:
    plt = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None

# --- Константы стиля ---
GH_GOLD = RGBColor(0xC5, 0x95, 0x00)
GH_DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
GH_BLACK = RGBColor(0x10, 0x10, 0x10)
GH_GRAY = RGBColor(0x80, 0x80, 0x80)

# Убедитесь, что этот путь к логотипу корректен
LOGO_PATH = 'app/static/img/logo_gh.png'


# === Вспомогательные функции ===

def _format_cell(cell, text, size=12, bold=False, align=PP_ALIGN.LEFT, color_rgb=None):
    if color_rgb is None:
        color_rgb = GH_BLACK

    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color_rgb
    p.alignment = align
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    cell.fill.background()


def _add_slide_footer(prs, slide):
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(15.0), Inches(8.3), height=Inches(0.6))
    except Exception:
        pass

    page_num_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.5), Inches(1.0), Inches(0.3))
    page_num_box.text_frame.text = f"Слайд {len(prs.slides)}"
    page_num_box.text_frame.paragraphs[0].font.size = Pt(10)
    page_num_box.text_frame.paragraphs[0].font.color.rgb = GH_GRAY


def _add_analyst_placeholder(slide):
    left = Inches(0.5)
    bottom = Inches(7.8)
    width = Inches(15.0)
    height = Inches(1.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, bottom, width, height
    )
    shape.line.color.rgb = GH_GOLD
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 248, 248)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), bottom, width - Inches(0.2), height
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Заключение ведущего аналитика:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = GH_DARK_BLUE

    p = tf.add_paragraph()
    p.text = "..."
    p.font.italic = True
    p.font.size = Pt(12)


def _create_bar_chart_image(labels, data, title):
    if plt is None:
        raise ImportError("Библиотека matplotlib не установлена.")

    fig, ax = plt.subplots(figsize=(14, 6))

    bars = ax.barh(labels, data, color=f"#{GH_GOLD.to_xml()[:6]}")

    ax.set_title(title, fontsize=16, fontweight='bold', color=f"#{GH_DARK_BLUE.to_xml()[:6]}")
    ax.tick_params(axis='y', labelsize=12, labelcolor=f"#{GH_BLACK.to_xml()[:6]}")
    ax.tick_params(axis='x', labelsize=12, labelcolor=f"#{GH_GRAY.to_xml()[:6]}")

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(f"#{GH_GRAY.to_xml()[:6]}")
    ax.spines['bottom'].set_color(f"#{GH_GRAY.to_xml()[:6]}")

    max_val = max(data) if data else 1
    for bar in bars:
        ax.text(bar.get_width() + (max_val * 0.01),
                bar.get_y() + bar.get_height() / 2,
                f'{bar.get_width():,.0f}',
                va='center', ha='left', fontsize=12, fontweight='bold',
                color=f"#{GH_DARK_BLUE.to_xml()[:6]}")

    fig.tight_layout()

    img_stream = io.BytesIO()
    fig.savefig(img_stream, format='png', dpi=100, transparent=True)
    plt.close(fig)
    img_stream.seek(0)
    return img_stream


# === Функции для создания слайдов ===

def _add_title_slide(prs, layout, complex_name):
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"Паспорт проекта: {complex_name or 'N/A'}"
    title.text_frame.paragraphs[0].font.color.rgb = GH_DARK_BLUE
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text = f"Аналитический дэшборд\nДата формирования: {date.today().strftime('%d.%m.%Y')}"

    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(0.5), height=Inches(1.0))
    except Exception as e:
        print(f"Warning: Не удалось добавить логотип {LOGO_PATH}: {e}")
        pass


def _add_kpi_slide(prs, layout, dynamic_data):
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.8))
    title_box.text_frame.text = "Ключевые показатели проекта"
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GH_DARK_BLUE

    kpis = [
        ("Общее кол-во квартир", dynamic_data.get('total_units', 0)),
        ("Остаток на сегодня", dynamic_data.get('total_remainders_count', 0)),
        ("Мес. до кадастра", dynamic_data.get('months_to_cadastre') or 'N/A'),
        ("Темп продаж (юнит/мес)", f"{dynamic_data.get('all_time_sales_pace', 0):.1f}")
    ]

    for i, (label, value) in enumerate(kpis):
        left = Inches(0.5 + i * 3.8)
        top = Inches(1.2)
        width = Inches(3.6)
        height = Inches(2.0)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = GH_GOLD
        shape.line.width = Pt(2.0)

        tf = shape.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = str(label)
        p.font.size = Pt(16)
        p.font.color.rgb = GH_GRAY
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = str(value)
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = GH_DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

    _add_analyst_placeholder(slide)
    _add_slide_footer(prs, slide)


def _add_team_leads_slide(prs, layout, static_data, dynamic_data):
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(7.5), Inches(0.8))
    title_box.text_frame.text = "Команда проекта"
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GH_DARK_BLUE

    team = [
        ("Руководитель проекта", static_data.get('project_manager') or '...'),
        ("Главный инженер", static_data.get('chief_engineer') or '...'),
        ("Руководитель ОП", static_data.get('sales_manager') or '...'),
    ]

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(7.5)
    height = Inches(4.0)

    tb_team = slide.shapes.add_table(len(team), 2, left, top, width, height).table
    tb_team.columns[0].width = Inches(3.0)
    tb_team.columns[1].width = Inches(4.5)

    for i, (role, name) in enumerate(team):
        _format_cell(tb_team.cell(i, 0), role, size=18, color_rgb=GH_GRAY)
        _format_cell(tb_team.cell(i, 1), name, size=18, bold=True, color_rgb=GH_BLACK)

    title_box_2 = slide.shapes.add_textbox(Inches(8.5), Inches(0.2), Inches(7.0), Inches(0.8))
    title_box_2.text_frame.text = "Лиды (Текущий месяц)"
    p = title_box_2.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GH_DARK_BLUE

    leads = dynamic_data.get('lead_stats', {})
    lead_data = [
        ("Всего заявок", leads.get('total_leads', 0)),
        ("Целевых заявок", leads.get('targeted_leads', 0)),
        ("Назначено встреч", leads.get('scheduled_meetings', 0)),
    ]

    left = Inches(8.5)
    width = Inches(7.0)

    tb_leads = slide.shapes.add_table(len(lead_data), 2, left, top, width, height).table
    tb_leads.columns[0].width = Inches(4.0)
    tb_leads.columns[1].width = Inches(3.0)

    for i, (label, value) in enumerate(lead_data):
        _format_cell(tb_leads.cell(i, 0), label, size=18, color_rgb=GH_GRAY)
        _format_cell(tb_leads.cell(i, 1), f"{value:,}", size=24, bold=True, align=PP_ALIGN.RIGHT,
                     color_rgb=GH_DARK_BLUE)

    _add_analyst_placeholder(slide)
    _add_slide_footer(prs, slide)


def _add_plan_fact_slide(prs, layout, dynamic_data, usd_rate):
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = "Показатели План-Факт (в USD)"
    title.text_frame.paragraphs[0].font.color.rgb = GH_DARK_BLUE

    rate = usd_rate if (usd_rate and usd_rate > 0) else 1.0

    lpf = dynamic_data.get('latest_plan_fact_data', {})
    td = dynamic_data.get('total_deviation_data', {})

    headers = ["Показатель", "План", "Факт", "Отклонение"]

    def get_row_data_usd(data_dict, key_plan, key_fact, is_currency=False):
        plan = data_dict.get(key_plan, 0) or 0
        fact = data_dict.get(key_fact, 0) or 0

        if is_currency:
            plan = plan / rate
            fact = fact / rate

        dev = fact - plan

        if is_currency:
            return f"${plan:,.0f}", f"${fact:,.0f}", f"${dev:,.0f}"
        else:
            return f"{plan:,.0f}", f"{fact:,.0f}", f"{dev:,.0f}"

    rows_lpf = [
        ("Продажи, шт.", *get_row_data_usd(lpf, 'plan_units', 'fact_units')),
        ("Контрактация, $", *get_row_data_usd(lpf, 'plan_volume', 'fact_volume', is_currency=True)),
        ("Поступления, $", *get_row_data_usd(lpf, 'plan_income', 'fact_income', is_currency=True)),
    ]

    rows_total = [
        ("Продажи, шт.", *get_row_data_usd(td, 'plan_units', 'fact_units')),
        ("Контрактация, $", *get_row_data_usd(td, 'plan_volume', 'fact_volume', is_currency=True)),
        ("Поступления, $", *get_row_data_usd(td, 'plan_income', 'fact_income', is_currency=True)),
    ]

    sub_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(0.5))
    sub_title.text_frame.text = f"Последний отчет ({lpf.get('period', 'N/A')})"
    sub_title.text_frame.paragraphs[0].font.size = Pt(20)

    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(15.0)
    height = Inches(2.5)

    table_lpf = slide.shapes.add_table(len(rows_lpf) + 1, 4, left, top, width, height).table
    table_lpf.columns[0].width = Inches(5.0)
    table_lpf.columns[1].width = Inches(3.3)
    table_lpf.columns[2].width = Inches(3.3)
    table_lpf.columns[3].width = Inches(3.3)

    for i, header in enumerate(headers):
        _format_cell(table_lpf.cell(0, i), header, bold=True, size=14, color_rgb=GH_DARK_BLUE,
                     align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for r, row_data in enumerate(rows_lpf, 1):
        _format_cell(table_lpf.cell(r, 0), row_data[0], size=14)
        _format_cell(table_lpf.cell(r, 1), row_data[1], size=14, align=PP_ALIGN.RIGHT)
        _format_cell(table_lpf.cell(r, 2), row_data[2], size=14, align=PP_ALIGN.RIGHT)

        deviation = float(row_data[3].replace(',', '').replace('$', ''))
        color = RGBColor(200, 0, 0) if deviation < 0 else RGBColor(0, 128, 0)
        _format_cell(table_lpf.cell(r, 3), row_data[3], size=14, bold=True, align=PP_ALIGN.RIGHT, color_rgb=color)

    sub_title_2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(15), Inches(0.5))
    sub_title_2.text_frame.text = "Суммарное отклонение (за все время)"
    sub_title_2.text_frame.paragraphs[0].font.size = Pt(20)

    top_2 = Inches(5.3)
    table_total = slide.shapes.add_table(len(rows_total) + 1, 4, left, top_2, width, height).table
    table_total.columns[0].width = Inches(5.0)
    table_total.columns[1].width = Inches(3.3)
    table_total.columns[2].width = Inches(3.3)
    table_total.columns[3].width = Inches(3.3)

    for i, header in enumerate(headers):
        _format_cell(table_total.cell(0, i), header, bold=True, size=14, color_rgb=GH_DARK_BLUE,
                     align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for r, row_data in enumerate(rows_total, 1):
        _format_cell(table_total.cell(r, 0), row_data[0], size=14)
        _format_cell(table_total.cell(r, 1), row_data[1], size=14, align=PP_ALIGN.RIGHT)
        _format_cell(table_total.cell(r, 2), row_data[2], size=14, align=PP_ALIGN.RIGHT)

        deviation = float(row_data[3].replace(',', '').replace('$', ''))
        color = RGBColor(200, 0, 0) if deviation < 0 else RGBColor(0, 128, 0)
        _format_cell(table_total.cell(r, 3), row_data[3], size=14, bold=True, align=PP_ALIGN.RIGHT, color_rgb=color)

    _add_analyst_placeholder(slide)
    _add_slide_footer(prs, slide)


def _add_remainders_slide(prs, layout, dynamic_data, usd_rate):
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = "Остатки по типам (на сегодня, в USD)"
    title.text_frame.paragraphs[0].font.color.rgb = GH_DARK_BLUE

    remainders = dynamic_data.get('remainders_by_type', {})
    rate = usd_rate if (usd_rate and usd_rate > 0) else 1.0

    if not remainders:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text_frame.text = "Данные по остаткам отсутствуют."
        _add_analyst_placeholder(slide)
        _add_slide_footer(prs, slide)
        return

    left = Inches(1.0)
    top = Inches(1.8)
    width = Inches(14.0)
    height = Inches(0.8) * (len(remainders) + 1)

    table = slide.shapes.add_table(len(remainders) + 1, 4, left, top, width, height).table
    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3.5)

    headers = ["Тип недвижимости", "Кол-во, шт.", "Общая стоимость (дно)", "Сред. цена дна ($/м²)"]
    for i, header in enumerate(headers):
        _format_cell(table.cell(0, i), header, bold=True, size=14, color_rgb=GH_DARK_BLUE,
                     align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    i = 1
    for prop_type, metrics in remainders.items():
        total_price_usd = (metrics.get('total_price', 0) or 0) / rate
        avg_price_usd = (metrics.get('avg_price_sqm', 0) or 0) / rate

        _format_cell(table.cell(i, 0), prop_type or "N/A", size=14)
        _format_cell(table.cell(i, 1), f"{metrics.get('count', 0):,}", size=14, align=PP_ALIGN.RIGHT)
        _format_cell(table.cell(i, 2), f"${total_price_usd:,.0f}", size=14, align=PP_ALIGN.RIGHT)
        _format_cell(table.cell(i, 3), f"${avg_price_usd:,.0f}", size=14, align=PP_ALIGN.RIGHT, bold=True,
                     color_rgb=GH_GOLD)
        i += 1

    _add_analyst_placeholder(slide)
    _add_slide_footer(prs, slide)


def _add_payment_distribution_slide(prs, layout, dynamic_data):
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.8))
    title_box.text_frame.text = "Структура продаж (по типу оплаты)"
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GH_DARK_BLUE

    payment_data = dynamic_data.get('payment_distribution', {})
    labels = payment_data.get('labels', [])
    data = payment_data.get('data', [])

    if not data or sum(data) == 0:
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(14.0), Inches(1.0))
        tb.text_frame.text = "Нет данных для построения графика."
        _add_analyst_placeholder(slide)
        _add_slide_footer(prs, slide)
        return

    try:
        chart_image_stream = _create_bar_chart_image(labels, data, "Количество сделок по типу оплаты")
        slide.shapes.add_picture(
            chart_image_stream,
            Inches(1.0), Inches(1.2),
            width=Inches(14.0)
        )
    except Exception as e:
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(14.0), Inches(1.0))
        tb.text_frame.text = f"Не удалось построить график: {e}"
        print(f"Ошибка Matplotlib: {e}")

    _add_analyst_placeholder(slide)
    _add_slide_footer(prs, slide)


def _add_competitors_slide(prs, layout, static_data, usd_rate):
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = "Анализ конкурентов"
    title.text_frame.paragraphs[0].font.color.rgb = GH_DARK_BLUE

    competitors = static_data.get('competitors', [])

    if not competitors:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text_frame.text = "Данные по конкурентам отсутствуют."
        _add_analyst_placeholder(slide)
        _add_slide_footer(prs, slide)
        return

    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(15.0)
    height = Inches(0.6) * (len(competitors) + 1)

    table = slide.shapes.add_table(len(competitors) + 1, 5, left, top, width, height).table

    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(2.5)
    table.columns[4].width = Inches(2.5)

    headers = ["Проект", "Класс", "План. сдача", "Темп, ю/мес", "Цена, $/м²"]
    for i, header in enumerate(headers):
        _format_cell(table.cell(0, i), header, bold=True, size=14, color_rgb=GH_DARK_BLUE,
                     align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    i = 1
    for comp in competitors:
        price_usd = (comp.get('price_per_sqm', 0) or 0)

        _format_cell(table.cell(i, 0), comp.get('competitor_name', 'N/A'), size=12, bold=True)
        _format_cell(table.cell(i, 1), comp.get('project_class', '-'), size=12, align=PP_ALIGN.CENTER)
        _format_cell(table.cell(i, 2), str(comp.get('planned_completion_date') or '-'), size=12, align=PP_ALIGN.CENTER)
        _format_cell(table.cell(i, 3), str(comp.get('sales_pace', '-') or '-'), size=12, align=PP_ALIGN.RIGHT)
        _format_cell(table.cell(i, 4), f"${price_usd:,.0f}" if price_usd > 0 else "-", size=14, align=PP_ALIGN.RIGHT,
                     bold=True, color_rgb=GH_GOLD)
        i += 1

    _add_analyst_placeholder(slide)
    _add_slide_footer(prs, slide)


# --- ГЛАВНАЯ ФУНКЦИЯ ГЕНЕРАЦИИ ---

def generate_passport_pptx(data: dict, usd_rate: float):
    """
    Генерирует PPTX презентацию, получая ГОТОВЫЕ данные.
    """
    if Presentation is None:
        raise ImportError("Библиотека python-pptx не установлена.")
    if plt is None:
        raise ImportError("Библиотека matplotlib не установлена.")

    static_data = data.get('static_data', {})
    dynamic_data = data.get('dynamic_data', {})
    complex_name = data.get('complex_name', 'N/A')

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    _add_title_slide(prs, title_slide_layout, complex_name)

    _add_kpi_slide(prs, blank_slide_layout, dynamic_data)

    _add_team_leads_slide(prs, blank_slide_layout, static_data, dynamic_data)

    _add_plan_fact_slide(prs, content_slide_layout, dynamic_data, usd_rate)

    _add_remainders_slide(prs, content_slide_layout, dynamic_data, usd_rate)

    _add_payment_distribution_slide(prs, blank_slide_layout, dynamic_data)

    _add_competitors_slide(prs, content_slide_layout, static_data, usd_rate)

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    return file_stream